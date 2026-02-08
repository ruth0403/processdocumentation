import streamlit as st
import pandas as pd
from pptx import Presentation
from transformers import pipeline

# ----------------------
# 1. Initialize LLM (HuggingFace free model)
# ----------------------
@st.cache_resource
def load_model():
    return pipeline("text2text-generation", model="meta-llama/Llama-2-7b-chat-hf")

question_gen = load_model()

# ----------------------
# 2. App Title
# ----------------------
st.title("📄 Free SOP Automation Portal")

# Initialize session state
if "submissions" not in st.session_state:
    st.session_state.submissions = []

if "questions" not in st.session_state:
    st.session_state.questions = []

# ----------------------
# 3. SOP Submission Form
# ----------------------
st.header("1️⃣ Submit SOP")
with st.form("sop_form"):
    process_name = st.text_input("Process Name")
    owner = st.text_input("Owner")
    steps = st.text_area("Steps (Separate each step with a newline)")
    inputs = st.text_input("Inputs")
    outputs = st.text_input("Outputs")
    submitted = st.form_submit_button("Submit SOP")

if submitted:
    sop = {
        "Process Name": process_name,
        "Owner": owner,
        "Steps": steps.split("\n"),
        "Inputs": inputs,
        "Outputs": outputs,
        "Status": "Pending",
        "Answers": ""
    }
    st.session_state.submissions.append(sop)
    st.success("SOP submitted successfully!")

# ----------------------
# 4. Display Submissions & Generate Questions
# ----------------------
st.header("2️⃣ Review Submissions & Clarify")

for idx, sub in enumerate(st.session_state.submissions):
    st.subheader(f"{idx+1}. {sub['Process Name']} (Status: {sub['Status']})")
    st.write(f"Owner: {sub['Owner']}")
    st.write(f"Steps: {sub['Steps']}")
    st.write(f"Inputs: {sub['Inputs']}")
    st.write(f"Outputs: {sub['Outputs']}")
    st.write(f"Clarifications/Answers: {sub['Answers']}")
    
    # Check for missing fields
    missing_fields = []
    for key in ["Process Name","Owner","Steps","Inputs","Outputs"]:
        if not sub[key] or (isinstance(sub[key], list) and len(sub[key]) == 0):
            missing_fields.append(key)
    if missing_fields:
        st.warning(f"Missing/unclear fields: {missing_fields}")

        # Generate clarifying questions using LLM
        prompt = f"Generate clarifying questions for the following SOP submission, focusing on missing fields {missing_fields}:\n{sub}"
        questions = question_gen(prompt, max_length=200)
        st.session_state.questions.append({"index": idx, "questions": questions[0]['generated_text']})
        st.info("**Clarifying Questions:**\n" + questions[0]['generated_text'])

# ----------------------
# 5. Employee Answers
# ----------------------
st.header("3️⃣ Answer Clarifying Questions")
for q in st.session_state.questions:
    idx = q["index"]
    st.subheader(f"Questions for SOP: {st.session_state.submissions[idx]['Process Name']}")
    st.write(q["questions"])
    answer = st.text_area(f"Your answers", key=f"answer_{idx}")
    if st.button(f"Submit Answers {idx}"):
        st.session_state.submissions[idx]["Answers"] = answer
        st.session_state.submissions[idx]["Status"] = "Clarified"
        st.success("Answers saved!")

# ----------------------
# 6. Generate PPT
# ----------------------
st.header("4️⃣ Generate SOP PPT")
if st.button("Generate PPT"):
    prs = Presentation()  # Or Presentation("template.pptx") if you have a template

    for sub in st.session_state.submissions:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content layout
        slide.shapes.title.text = sub["Process Name"]

        content_box = slide.placeholders[1]
        content_text = f"""
Owner: {sub['Owner']}

Steps:
"""
        for step in sub["Steps"]:
            content_text += f"• {step}\n"

        content_text += f"""
Inputs: {sub['Inputs']}
Outputs: {sub['Outputs']}
Clarifications: {sub.get('Answers','N/A')}
"""
        content_box.text = content_text

    prs.save("SOP_Documented.pptx")
    st.success("PPT generated successfully!")

    with open("SOP_Documented.pptx", "rb") as f:
        st.download_button("Download SOP PPT", f.read(), "SOP_Documented.pptx")
